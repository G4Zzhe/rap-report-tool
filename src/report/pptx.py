"""PPT 报告生成。

读取公司 PPT 模板，基于数据、AI 文案和图表自动填充以下内容：
- 封面（标题、副标题、周期）
- 数据概览
- 平台榜单总结 + 各平台 TOP10
- 头部艺人表现
- 跨平台爆款
- 可视化图表页
- 待人工补充的章节提示页

其余页面保留模板原有内容，供人工后续补充。
"""
import logging
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.utils import OUTPUT_DIR

logger = logging.getLogger("rap_report")

# 模板路径（Windows 路径，WSL 下映射为 /mnt/d/...）
DEFAULT_TEMPLATE_PATH = Path("/mnt/d/G4Zzz/工作文档/神速信息_说唱音乐行业全景洞察分析报告（2026.7.16-7.31）.pptx")


def _format_period(start_date: str, end_date: str) -> str:
    """格式化为中文周期字符串。"""
    try:
        start = datetime.strptime(start_date, "%Y-%m-%d")
        end = datetime.strptime(end_date, "%Y-%m-%d")
        return f"报告周期：{start.year}年{start.month}月{start.day}日 - {end.month}月{end.day}日"
    except ValueError:
        return f"报告周期：{start_date} - {end_date}"


def _find_shape_by_text(slide, keyword: str) -> Optional[Any]:
    """在幻灯片中查找包含指定关键词的第一个文本形状。"""
    for shape in slide.shapes:
        if shape.has_text_frame and keyword in shape.text_frame.text:
            return shape
    return None


def _set_text(shape: Any, text: str) -> None:
    """设置形状文本，保留原有格式。"""
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.text = text


def _get_blank_layout(prs: Presentation):
    """获取空白布局，兼容不同模板。"""
    for layout in prs.slide_layouts:
        if "空白" in layout.name or "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]


def _add_title_slide(prs: Presentation, title: str, subtitle: str, period: str) -> None:
    """添加封面页。"""
    layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.2), Inches(8.0), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 153)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    sub_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(4.0), Inches(8.0), Inches(1.0)
    )
    tf = sub_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.CENTER

    # 周期
    period_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(5.2), Inches(8.0), Inches(0.6)
    )
    tf = period_box.text_frame
    tf.text = period
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER


def _add_overview_slide(prs: Presentation, analysis: Dict[str, Any], start_date: str, end_date: str) -> None:
    """添加数据概览页。"""
    layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9.0), Inches(0.8)
    )
    tf = title_box.text_frame
    tf.text = "数据概览"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 153)

    stats = [
        ("本期采集歌曲", f"{analysis['total_songs']} 首"),
        ("涉及平台", f"{analysis['total_platforms']} 个"),
        ("头部艺人数量", f"{len(analysis['artist_summary'])}"),
        ("跨平台爆款", f"{len(analysis['cross_platform_hits'])} 首"),
        ("报告周期", f"{start_date} 至 {end_date}"),
    ]

    top = 1.6
    for label, value in stats:
        box = slide.shapes.add_textbox(Inches(1.0), Inches(top), Inches(8.0), Inches(0.6))
        tf = box.text_frame
        tf.text = f"{label}：{value}"
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(50, 50, 50)
        top += 0.8


def _add_platform_summary_slide(
    prs: Presentation,
    analysis: Dict[str, Any],
    ai_texts: Dict[str, str],
) -> None:
    """添加平台榜单总结页。"""
    layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9.0), Inches(0.8)
    )
    tf = title_box.text_frame
    tf.text = "平台榜单总结"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 153)

    content = ai_texts.get("platform_summary") or "（AI 文案未生成）"
    content_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5), Inches(8.4), Inches(2.5)
    )
    tf = content_box.text_frame
    tf.text = content
    tf.word_wrap = True
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(60, 60, 60)


def _add_top10_slides(prs: Presentation, analysis: Dict[str, Any]) -> None:
    """为每个榜单添加 TOP10 页。"""
    raw = analysis.get("raw")
    if raw is None or raw.empty:
        return

    for chart_name, group in raw.groupby("chart_name"):
        layout = _get_blank_layout(prs)
        slide = prs.slides.add_slide(layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.7)
        )
        tf = title_box.text_frame
        tf.text = f"{chart_name} TOP10"
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 78, 153)

        top10 = group.nsmallest(10, "rank")[["rank", "song_name", "artist"]]
        rows = len(top10)
        cols = 3
        left = Inches(0.8)
        top = Inches(1.3)
        width = Inches(8.4)
        height = Inches(0.5 * rows + 0.2)

        table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

        # 表头
        headers = ["排名", "歌曲", "艺人"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(31, 78, 153)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        for idx, (_, row) in enumerate(top10.iterrows(), start=1):
            table.cell(idx, 0).text = str(int(row["rank"]))
            table.cell(idx, 1).text = str(row["song_name"])
            table.cell(idx, 2).text = str(row["artist"])
            for col in range(cols):
                cell = table.cell(idx, col)
                cell.text_frame.paragraphs[0].font.size = Pt(12)


def _add_artist_slide(prs: Presentation, analysis: Dict[str, Any], ai_texts: Dict[str, str]) -> None:
    """添加头部艺人表现页。"""
    layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9.0), Inches(0.8)
    )
    tf = title_box.text_frame
    tf.text = "头部艺人表现"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 153)

    content = ai_texts.get("artist_insight") or "（AI 文案未生成）"
    content_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.4), Inches(8.4), Inches(1.8)
    )
    tf = content_box.text_frame
    tf.text = content
    tf.word_wrap = True
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(60, 60, 60)

    artist_df = analysis["artist_summary"].head(10)
    if artist_df.empty:
        return

    rows = len(artist_df)
    cols = 5
    left = Inches(0.5)
    top = Inches(3.4)
    width = Inches(9.0)
    height = Inches(0.45 * rows + 0.2)

    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
    headers = ["艺人", "总上榜数", "平台数", "最佳排名", "最佳歌曲"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(31, 78, 153)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    for idx, (_, row) in enumerate(artist_df.iterrows(), start=1):
        table.cell(idx, 0).text = str(row["artist"])
        table.cell(idx, 1).text = str(int(row["total_entries"]))
        table.cell(idx, 2).text = str(int(row["platform_count"]))
        table.cell(idx, 3).text = str(int(row["best_rank"]))
        table.cell(idx, 4).text = str(row["best_song"])
        for col in range(cols):
            table.cell(idx, col).text_frame.paragraphs[0].font.size = Pt(11)


def _add_hit_songs_slide(prs: Presentation, analysis: Dict[str, Any], ai_texts: Dict[str, str]) -> None:
    """添加跨平台爆款页。"""
    layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9.0), Inches(0.8)
    )
    tf = title_box.text_frame
    tf.text = "跨平台爆款歌曲"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 153)

    content = ai_texts.get("hit_songs_insight") or "（AI 文案未生成）"
    content_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.4), Inches(8.4), Inches(1.5)
    )
    tf = content_box.text_frame
    tf.text = content
    tf.word_wrap = True
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(60, 60, 60)

    hits_df = analysis["cross_platform_hits"]
    if hits_df.empty:
        info = slide.shapes.add_textbox(
            Inches(0.8), Inches(3.2), Inches(8.4), Inches(0.8)
        )
        info.text_frame.text = "本期暂无跨平台爆款歌曲。"
        info.text_frame.paragraphs[0].font.size = Pt(18)
        return

    rows = len(hits_df)
    cols = 4
    left = Inches(1.0)
    top = Inches(3.2)
    width = Inches(8.0)
    height = Inches(0.5 * rows + 0.2)

    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
    headers = ["歌曲", "艺人", "上榜平台数", "最佳排名"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(31, 78, 153)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    for idx, (_, row) in enumerate(hits_df.iterrows(), start=1):
        table.cell(idx, 0).text = str(row["song_name"])
        table.cell(idx, 1).text = str(row["artist"])
        table.cell(idx, 2).text = str(int(row["platform_count"]))
        table.cell(idx, 3).text = str(int(row["best_rank"]))
        for col in range(cols):
            table.cell(idx, col).text_frame.paragraphs[0].font.size = Pt(12)


def _add_chart_slides(prs: Presentation, chart_paths: List[Path]) -> None:
    """为每个图表添加一页。"""
    for chart_path in chart_paths:
        if not chart_path.exists():
            continue
        layout = _get_blank_layout(prs)
        slide = prs.slides.add_slide(layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.7)
        )
        tf = title_box.text_frame
        tf.text = chart_path.stem.replace("chart_", "").replace("_", " ").title()
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 78, 153)

        slide.shapes.add_picture(
            str(chart_path),
            Inches(1.0),
            Inches(1.3),
            width=Inches(8.0),
        )


def _add_placeholder_slide(prs: Presentation, section: str, items: List[str]) -> None:
    """添加待人工补充的章节提示页。"""
    layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9.0), Inches(0.8)
    )
    tf = title_box.text_frame
    tf.text = f"{section}（待补充）"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 80, 80)

    top = 1.6
    for item in items:
        box = slide.shapes.add_textbox(
            Inches(1.0), Inches(top), Inches(8.0), Inches(0.5)
        )
        tf = box.text_frame
        tf.text = f"• {item}"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(80, 80, 80)
        top += 0.6


def generate_pptx(
    analysis: Dict[str, Any],
    ai_texts: Dict[str, str],
    chart_paths: List[Path],
    start_date: str,
    end_date: str,
    config: Dict[str, Any],
    template_path: Optional[Path] = None,
) -> Path:
    """生成 PPT 报告初稿。

    优先使用公司模板；若模板不可用，则创建一个包含基础页面的新演示文稿。
    """
    report_cfg = config.get("report", {})
    title = report_cfg.get("title", "说唱音乐行业全景洞察分析报告")
    subtitle = report_cfg.get("subtitle", "深度解析舆情、数据、艺人及产业链趋势（双周版）")
    period = _format_period(start_date, end_date)

    template = template_path or DEFAULT_TEMPLATE_PATH
    if template and template.exists():
        prs = Presentation(str(template))
        logger.info("使用 PPT 模板: %s", template)
    else:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        logger.warning("PPT 模板未找到，使用空白演示文稿: %s", template)

    # 删除模板中已有的占位内容页（可选），保留原有结构，追加自动化页面
    _add_title_slide(prs, title, subtitle, period)
    _add_overview_slide(prs, analysis, start_date, end_date)
    _add_platform_summary_slide(prs, analysis, ai_texts)
    _add_top10_slides(prs, analysis)
    _add_artist_slide(prs, analysis, ai_texts)
    _add_hit_songs_slide(prs, analysis, ai_texts)
    _add_chart_slides(prs, chart_paths)

    # 待人工补充的章节提示
    _add_placeholder_slide(
        prs,
        "舆情与话题",
        ["全网舆情与话题总览", "正面/中性/负面舆情关注点"],
    )
    _add_placeholder_slide(
        prs,
        "厂牌与行业",
        ["代表厂牌深度解析", "周期内关键行业信息", "艺人动态与商业变现"],
    )
    _add_placeholder_slide(
        prs,
        "产业与风险",
        ["产业链与资本动态", "政策监管与合规", "风险机遇与策略建议"],
    )

    filename = f"report_{start_date.replace('-', '')}_{end_date.replace('-', '')}.pptx"
    path = OUTPUT_DIR / filename
    prs.save(str(path))
    logger.info("PPT 报告已保存: %s", path)
    return path
